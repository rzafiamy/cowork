"""
📄 Document Creation Tools
Tools for generating PDF, PPTX (PowerPoint), XLSX (Excel), and DOCX (Word)
files directly into the session workspace artifacts folder.
All file I/O is routed through file_manager (ACL-enforced).

Libraries:
  • PDF  → reportlab    (vector PDF, no external deps)
  • PPTX → python-pptx  (native .pptx files)
  • XLSX → openpyxl     (native .xlsx files)
  • DOCX → python-docx  (native .docx files)
"""

import json
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional

from ..base import BaseTool
from ...workspace import workspace_manager, WORKSPACE_ROOT
from ...acl import file_manager


# ─── Workspace helper ─────────────────────────────────────────────────────────

def _get_artifacts_dir(scratchpad) -> Path:
    """Return the workspace artifacts/ path, falling back to WORKSPACE_ROOT."""
    if scratchpad:
        ws = workspace_manager.get_by_session_id(scratchpad.session_id)
        if ws:
            return ws.artifacts_path
    fallback = WORKSPACE_ROOT / "artifacts"
    fallback.mkdir(parents=True, exist_ok=True)
    return fallback


def _safe_filename(name: str) -> str:
    """Strip path traversal, keep only the base name."""
    return Path(name).name


def _workspace_rel_path(path: Path) -> str:
    """Best-effort workspace-relative path for user-facing output."""
    if path.parent.name == "artifacts":
        return f"artifacts/{path.name}"
    return path.name


# ─── PDF ─────────────────────────────────────────────────────────────────────

class DocumentCreatePdfTool(BaseTool):

    @property
    def name(self) -> str:
        return "document_create_pdf"

    @property
    def description(self) -> str:
        return (
            "Create a PDF document and save it to the workspace artifacts folder. "
            "Accepts a title and a list of sections as JSON: "
            "[{\"heading\": \"...\", \"text\": \"...\", \"bullets\": [\"...\"]}, ...]. "
            "Returns the absolute file path for use with email attachments."
        )

    @property
    def category(self) -> str:
        return "DOCUMENT_TOOLS"

    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Output filename, e.g. 'report.pdf'",
                },
                "title": {
                    "type": "string",
                    "description": "Document title shown at the top",
                },
                "sections": {
                    "type": "string",
                    "description": (
                        "JSON array of sections. Each section: "
                        "{\"heading\": str (optional), \"text\": str (optional), "
                        "\"bullets\": [str] (optional), \"image\": str (absolute path, optional)}. "
                        "Example: [{\"heading\": \"Intro\", \"text\": \"Hello world.\", \"bullets\": [\"Point A\", \"Point B\"], \"image\": \"/path/to/img.png\"}]"
                    ),
                },
                "author": {
                    "type": "string",
                    "description": "Author name shown in the document footer (optional)",
                },
            },
            "required": ["filename", "title", "sections"],
        }

    def execute(self, filename: str, title: str, sections: str, author: str = "") -> str:
        self._emit(f"📄 Creating PDF: '{filename}'...")
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
            from reportlab.lib import colors
            from reportlab.platypus import (
                SimpleDocTemplate, Paragraph, Spacer,
                ListFlowable, ListItem, HRFlowable,
            )
        except ImportError:
            return "❌ reportlab not installed. Run: pip install reportlab"

        try:
            section_data: List[dict] = json.loads(sections)
        except json.JSONDecodeError as e:
            return f"❌ Invalid JSON in 'sections': {e}"

        artifacts_dir = _get_artifacts_dir(self.scratchpad)
        out_path = artifacts_dir / _safe_filename(filename)
        if not out_path.suffix:
            out_path = out_path.with_suffix(".pdf")

        buf = BytesIO()
        doc = SimpleDocTemplate(
            buf,
            pagesize=A4,
            rightMargin=20 * mm,
            leftMargin=20 * mm,
            topMargin=25 * mm,
            bottomMargin=20 * mm,
            title=title,
            author=author,
        )

        styles = getSampleStyleSheet()
        # Custom styles
        title_style = ParagraphStyle(
            "DocTitle",
            parent=styles["Title"],
            fontSize=24,
            spaceAfter=8,
            textColor=colors.HexColor("#1a1a2e"),
        )
        h1_style = ParagraphStyle(
            "H1",
            parent=styles["Heading1"],
            fontSize=15,
            spaceBefore=14,
            spaceAfter=4,
            textColor=colors.HexColor("#16213e"),
        )
        body_style = ParagraphStyle(
            "Body",
            parent=styles["Normal"],
            fontSize=10,
            leading=15,
            spaceAfter=6,
        )
        bullet_style = ParagraphStyle(
            "Bullet",
            parent=styles["Normal"],
            fontSize=10,
            leading=14,
            leftIndent=12,
        )

        story = [
            Paragraph(title, title_style),
            HRFlowable(width="100%", thickness=1, color=colors.HexColor("#16213e")),
            Spacer(1, 6 * mm),
        ]
        if author:
            story.append(Paragraph(f"<i>Author: {author}</i>", styles["Italic"]))
            story.append(Spacer(1, 4 * mm))

        for sec in section_data:
            heading = sec.get("heading", "")
            text = sec.get("text", "")
            bullets = sec.get("bullets", [])

            if heading:
                story.append(Paragraph(heading, h1_style))
            if text:
                # Escape XML special chars for ReportLab
                safe_text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(safe_text, body_style))
            if bullets:
                items = [
                    ListItem(Paragraph(b.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), bullet_style))
                    for b in bullets
                ]
                story.append(ListFlowable(items, bulletType="bullet", start="•", leftIndent=18))
            
            image_path = sec.get("image")
            if image_path and Path(image_path).exists():
                from reportlab.platypus import Image
                try:
                    img = Image(image_path)
                    max_width = 150 * mm
                    if img.drawWidth > max_width:
                        ratio = max_width / img.drawWidth
                        img.drawWidth = max_width
                        img.drawHeight = img.drawHeight * ratio
                    story.append(img)
                except Exception as e:
                    story.append(Paragraph(f"[Image Error: {e}]", body_style))

            story.append(Spacer(1, 3 * mm))

        doc.build(story)
        file_manager.write_bytes(out_path, buf.getvalue(), reason=f"document_create_pdf: {filename}")
        size_kb = out_path.stat().st_size // 1024
        rel_path = _workspace_rel_path(out_path)

        return (
            f"✅ PDF created!\n"
            f"• File: `{out_path.name}`\n"
            f"• Workspace path: `{rel_path}`\n"
            f"• Path: `{out_path}`\n"
            f"• Open: `/open {rel_path}`\n"
            f"• Size: {size_kb} KB\n"
            f"• Sections: {len(section_data)}"
        )


# ─── PPTX ────────────────────────────────────────────────────────────────────

class DocumentCreatePptxTool(BaseTool):

    @property
    def name(self) -> str:
        return "document_create_pptx"

    @property
    def description(self) -> str:
        return (
            "Create a PowerPoint (.pptx) presentation and save it to the workspace artifacts folder. "
            "Accepts a list of slides as JSON: "
            "[{\"title\": \"...\", \"content\": \"...\", \"bullets\": [\"...\"], \"layout\": \"auto|full_text|centered|section|text_image\"}, ...]. "
            "The first slide is automatically a title slide. "
            "Returns the absolute file path for use with email attachments."
        )

    @property
    def category(self) -> str:
        return "DOCUMENT_TOOLS"

    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Output filename, e.g. 'presentation.pptx'",
                },
                "title": {
                    "type": "string",
                    "description": "Presentation overall title (used on the cover slide)",
                },
                "subtitle": {
                    "type": "string",
                    "description": "Subtitle or author shown on the cover slide (optional)",
                },
                "slides": {
                    "type": "string",
                    "description": (
                        "JSON array of slides. Each slide: "
                        "{\"title\": str, \"content\": str (optional), \"bullets\": [str] (optional), "
                        "\"image\": str (absolute path, optional), \"key_message\": str (optional), "
                        "\"layout\": \"auto|full_text|centered|section|text_image\" (optional)}. "
                        "Example: [{\"title\": \"Introduction\", \"bullets\": [\"Key point 1\"], \"layout\": \"auto\"}]"
                    ),
                },
                "theme_color": {
                    "type": "string",
                    "description": "Hex color for the accent/header color, e.g. '#2563EB' (optional, default dark blue)",
                },
                "design_preset": {
                    "type": "string",
                    "description": "Visual style preset: 'executive' (default), 'bold', 'minimal', or 'modern-light'",
                },
            },
            "required": ["filename", "title", "slides"],
        }

    def execute(
        self,
        filename: str,
        title: str,
        slides: str,
        subtitle: str = "",
        theme_color: str = "#1a1a2e",
        design_preset: str = "executive",
    ) -> str:
        self._emit(f"📊 Creating PPTX: '{filename}'...")
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return "❌ python-pptx not installed. Run: pip install python-pptx"

        try:
            slides_data: List[dict] = json.loads(slides)
        except json.JSONDecodeError as e:
            return f"❌ Invalid JSON in 'slides': {e}"

        artifacts_dir = _get_artifacts_dir(self.scratchpad)
        out_path = artifacts_dir / _safe_filename(filename)
        if not out_path.suffix:
            out_path = out_path.with_suffix(".pptx")

        # Parse hex color to RGB tuple.
        def _rgb_tuple(hex_color: str, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
            h = hex_color.lstrip("#")
            if len(h) == 3:
                h = "".join(c * 2 for c in h)
            try:
                return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            except Exception:
                return fallback

        def _as_rgb(value: tuple[int, int, int]) -> RGBColor:
            return RGBColor(value[0], value[1], value[2])

        def _mix(c1: tuple[int, int, int], c2: tuple[int, int, int], ratio: float) -> tuple[int, int, int]:
            return (
                int(c1[0] * (1 - ratio) + c2[0] * ratio),
                int(c1[1] * (1 - ratio) + c2[1] * ratio),
                int(c1[2] * (1 - ratio) + c2[2] * ratio),
            )

        preset = (design_preset or "executive").strip().lower()
        if preset not in {"executive", "bold", "minimal", "modern-light"}:
            preset = "executive"

        palettes = {
            "executive": {
                "accent": (22, 78, 99),
                "bg": (244, 247, 250),
                "surface": (230, 238, 245),
                "text": (20, 30, 46),
                "muted": (106, 121, 141),
            },
            "bold": {
                "accent": (190, 52, 85),
                "bg": (250, 242, 244),
                "surface": (243, 222, 228),
                "text": (45, 20, 30),
                "muted": (126, 80, 95),
            },
            "minimal": {
                "accent": (52, 98, 138),
                "bg": (249, 250, 251),
                "surface": (236, 240, 245),
                "text": (30, 37, 46),
                "muted": (112, 122, 132),
            },
            "modern-light": {
                "accent": (59, 130, 246),
                "bg": (255, 255, 255),
                "surface": (248, 250, 252),
                "text": (30, 41, 59),
                "muted": (100, 116, 139),
            },
        }
        palette = palettes[preset].copy()
        palette["accent"] = _rgb_tuple(theme_color, palette["accent"])
        palette["accent_dark"] = _mix(palette["accent"], (0, 0, 0), 0.28)
        palette["accent_soft"] = _mix(palette["accent"], (255, 255, 255), 0.75)

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # blank

        def _add_rect(slide, left, top, width, height, fill_rgb):
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _as_rgb(fill_rgb)
            shape.line.fill.background()
            return shape

        def _add_textbox(
            slide,
            left,
            top,
            width,
            height,
            text,
            font_size,
            bold=False,
            color=None,
            align=PP_ALIGN.LEFT,
            word_wrap=True,
            font_name="Aptos",
        ):
            txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txb.text_frame
            tf.word_wrap = word_wrap
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = font_name
            if color:
                run.font.color.rgb = _as_rgb(color)
            return txb

        def _trim_text(value: str, max_chars: int) -> str:
            s = str(value or "").strip()
            if len(s) <= max_chars:
                return s
            return s[: max_chars - 3].rstrip() + "..."

        # Cover slide
        cover = prs.slides.add_slide(blank_layout)
        _add_rect(cover, 0, 0, 13.33, 7.5, palette["bg"])
        if preset == "modern-light":
            # Subtle accent line instead of heavy blocks
            _add_rect(cover, 0, 0, 13.33, 0.15, palette["accent"])
        else:
            _add_rect(cover, 0, 0, 13.33, 1.2, palette["accent_dark"])
            _add_rect(cover, 0, 5.9, 13.33, 1.6, palette["accent"])
            _add_rect(cover, 8.8, 1.2, 4.53, 4.7, palette["accent_soft"])
        _add_textbox(
            cover,
            0.8,
            1.7,
            7.6,
            2.2,
            _trim_text(title, 90),
            42,
            bold=True,
            color=palette["text"],
            align=PP_ALIGN.LEFT,
            font_name="Aptos Display",
        )
        if subtitle:
            _add_textbox(
                cover,
                0.8,
                4.25,
                11.7,
                1.0,
                _trim_text(subtitle, 120),
                20,
                color=palette["muted"],
                align=PP_ALIGN.LEFT,
                font_name="Aptos",
            )
        _add_textbox(
            cover,
            0.8,
            6.25,
            12.0,
            0.6,
            "Built with Cowork Deck Studio",
            11,
            color=(245, 247, 252),
            align=PP_ALIGN.LEFT,
            font_name="Aptos",
        )

        # Content slides with alternating layouts and message callouts.
        for i, slide_data in enumerate(slides_data):
            sl = prs.slides.add_slide(blank_layout)
            slide_title = _trim_text(slide_data.get("title", f"Slide {i + 1}"), 85)
            content_text = _trim_text(slide_data.get("content", ""), 280)
            bullets = [str(b).strip() for b in slide_data.get("bullets", []) if str(b).strip()][:6]
            bullets = [_trim_text(b, 92) for b in bullets]
            key_message = _trim_text(slide_data.get("key_message", ""), 140)
            layout = (slide_data.get("layout", "auto") or "auto").strip().lower()

            image_path = slide_data.get("image")
            has_image = bool(image_path and Path(image_path).exists())

            if layout not in {"auto", "section", "text_image", "centered", "full_text"}: # Updated valid layouts
                layout = "auto"
            if layout == "auto":
                if not content_text and not bullets:
                    layout = "centered"
                elif has_image:
                    layout = "text_image"
                elif len(content_text or "") < 100 and len(bullets or []) < 2:
                    layout = "centered"
                else:
                    layout = "full_text"

            _add_rect(sl, 0, 0, 13.33, 7.5, palette["bg"])
            if preset == "modern-light":
                _add_rect(sl, 0, 0, 13.33, 0.08, palette["accent"])
            else:
                _add_rect(sl, 0, 0, 13.33, 0.85, palette["accent_dark"])

            title_color = palette["accent_dark"] if preset == "modern-light" else (250, 252, 255)
            _add_textbox(
                sl,
                0.35,
                0.11,
                11.6,
                0.55,
                slide_title,
                21,
                bold=True,
                color=title_color,
                align=PP_ALIGN.LEFT,
                font_name="Aptos Display",
            )
            page_num_color = palette["muted"] if preset == "modern-light" else (224, 232, 242)
            _add_textbox(
                sl,
                12.2,
                0.12,
                0.95,
                0.45,
                str(i + 2),
                12,
                color=page_num_color,
                align=PP_ALIGN.RIGHT,
                font_name="Aptos",
            )

            if layout == "section":
                _add_rect(sl, 0, 1.1, 13.33, 6.4, palette["accent_soft"])
                _add_textbox(
                    sl,
                    1.0,
                    2.2,
                    11.3,
                    1.6,
                    slide_title,
                    44,
                    bold=True,
                    color=palette["accent_dark"],
                    align=PP_ALIGN.CENTER,
                    font_name="Aptos Display",
                )
                if content_text:
                    _add_textbox(
                        sl,
                        1.2,
                        4.2,
                        10.9,
                        1.3,
                        content_text,
                        18,
                        color=palette["text"],
                        align=PP_ALIGN.CENTER,
                        font_name="Aptos",
                    )
                continue

            if layout == "centered":
                _add_textbox(
                    sl,
                    1.5,
                    2.8,
                    10.33,
                    1.5,
                    slide_title,
                    36,
                    bold=True,
                    color=palette["accent"],
                    align=PP_ALIGN.CENTER,
                    font_name="Aptos Display",
                )
                if content_text:
                    _add_textbox(
                        sl,
                        2.0,
                        4.5,
                        9.33,
                        1.2,
                        content_text,
                        20,
                        color=palette["text"],
                        align=PP_ALIGN.CENTER,
                        font_name="Aptos",
                    )
                continue

            if layout == "full_text":
                y_offset = 1.3
                if content_text:
                    _add_textbox(
                        sl, 0.8, y_offset, 11.7, 1.2, content_text, 18,
                        color=palette["text"], align=PP_ALIGN.LEFT
                    )
                    y_offset += 1.4
                if bullets:
                    bullet_box = sl.shapes.add_textbox(Inches(0.8), Inches(y_offset), Inches(11.7), Inches(4.5))
                    tf = bullet_box.text_frame
                    tf.word_wrap = True
                    for idx, bullet in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        p.text = f"• {bullet}"
                        p.space_before = Pt(8)
                        for run in p.runs:
                            run.font.size = Pt(20)
                            run.font.name = "Aptos"
                            run.font.color.rgb = _as_rgb(palette["text"])
                continue

            # Default: text_image (Split Layout)
            image_on_right = (i % 2 == 0)
            text_left = 0.6 if image_on_right else 6.95
            image_left = 6.95 if image_on_right else 0.6

            _add_rect(sl, text_left - 0.15, 1.05, 5.8, 5.7, palette["surface"])
            _add_rect(sl, image_left - 0.15, 1.05, 5.8, 5.7, _mix(palette["accent_soft"], palette["bg"], 0.45))

            y_offset = 1.35
            if content_text:
                _add_textbox(
                    sl,
                    text_left + 0.15,
                    y_offset,
                    5.2,
                    1.15,
                    content_text,
                    17,
                    color=palette["text"],
                    align=PP_ALIGN.LEFT,
                    font_name="Aptos",
                )
                y_offset += 1.25

            if bullets:
                bullet_box = sl.shapes.add_textbox(Inches(text_left + 0.15), Inches(y_offset), Inches(5.2), Inches(3.2))
                tf = bullet_box.text_frame
                tf.word_wrap = True
                for idx, bullet in enumerate(bullets):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.text = f"• {bullet}"
                    p.space_before = Pt(6)
                    for run in p.runs:
                        run.font.size = Pt(16)
                        run.font.name = "Aptos"
                        run.font.color.rgb = _as_rgb(palette["text"])

            if has_image:
                try:
                    sl.shapes.add_picture(str(image_path), Inches(image_left + 0.08), Inches(1.18), width=Inches(5.55), height=Inches(4.7))
                except Exception:
                    _add_textbox(
                        sl,
                        image_left + 0.55,
                        3.0,
                        4.7,
                        0.8,
                        "Image could not be rendered",
                        14,
                        color=palette["muted"],
                        align=PP_ALIGN.CENTER,
                        font_name="Aptos",
                    )
            else:
                _add_textbox(
                    sl,
                    image_left + 0.6,
                    3.0,
                    4.7,
                    0.8,
                    "Add chart, screenshot, or visual proof here",
                    14,
                    color=palette["muted"],
                    align=PP_ALIGN.CENTER,
                    font_name="Aptos",
                )

            if key_message:
                _add_rect(sl, 0.55, 6.35, 12.2, 0.78, palette["accent"])
                _add_textbox(
                    sl,
                    0.85,
                    6.5,
                    11.7,
                    0.45,
                    f"Key message: {key_message}",
                    13,
                    bold=True,
                    color=(250, 252, 255),
                    align=PP_ALIGN.LEFT,
                    font_name="Aptos",
                )

        prs.save(str(out_path))
        size_kb = out_path.stat().st_size // 1024
        rel_path = _workspace_rel_path(out_path)

        return (
            f"✅ PPTX created!\n"
            f"• File: `{out_path.name}`\n"
            f"• Workspace path: `{rel_path}`\n"
            f"• Path: `{out_path}`\n"
            f"• Open: `/open {rel_path}`\n"
            f"• Style: {preset}\n"
            f"• Size: {size_kb} KB\n"
            f"• Slides: {1 + len(slides_data)} (1 cover + {len(slides_data)} content)"
        )


# ─── XLSX ────────────────────────────────────────────────────────────────────

class DocumentCreateXlsxTool(BaseTool):

    @property
    def name(self) -> str:
        return "document_create_xlsx"

    @property
    def description(self) -> str:
        return (
            "Create an Excel (.xlsx) spreadsheet and save it to the workspace artifacts folder. "
            "Accepts sheet data as JSON: "
            "{\"Sheet Name\": {\"headers\": [\"col1\", \"col2\"], \"rows\": [[\"a\", 1], [\"b\", 2]]}, ...}. "
            "Supports multiple sheets, auto-column widths, and styled headers. "
            "Returns the absolute file path."
        )

    @property
    def category(self) -> str:
        return "DOCUMENT_TOOLS"

    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Output filename, e.g. 'data.xlsx'",
                },
                "sheets": {
                    "type": "string",
                    "description": (
                        "JSON object mapping sheet name → sheet data. "
                        "Sheet data: {\"headers\": [str, ...], \"rows\": [[value, ...], ...]}. "
                        "Example: {\"Sales\": {\"headers\": [\"Month\", \"Revenue\"], \"rows\": [[\"Jan\", 50000], [\"Feb\", 62000]]}}"
                    ),
                },
                "title": {
                    "type": "string",
                    "description": "Document title stored in file metadata (optional)",
                },
            },
            "required": ["filename", "sheets"],
        }

    def execute(self, filename: str, sheets: str, title: str = "") -> str:
        self._emit(f"📊 Creating XLSX: '{filename}'...")
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
        except ImportError:
            return "❌ openpyxl not installed. Run: pip install openpyxl"

        try:
            sheets_data: dict = json.loads(sheets)
        except json.JSONDecodeError as e:
            return f"❌ Invalid JSON in 'sheets': {e}"

        artifacts_dir = _get_artifacts_dir(self.scratchpad)
        out_path = artifacts_dir / _safe_filename(filename)
        if not out_path.suffix:
            out_path = out_path.with_suffix(".xlsx")

        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # Remove default sheet

        # Styles
        header_font = Font(bold=True, color="FFFFFF", size=11)
        header_fill = PatternFill("solid", fgColor="1A1A2E")
        center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        left_align = Alignment(horizontal="left", vertical="center", wrap_text=True)
        thin_border = Border(
            left=Side(style="thin", color="DDDDDD"),
            right=Side(style="thin", color="DDDDDD"),
            top=Side(style="thin", color="DDDDDD"),
            bottom=Side(style="thin", color="DDDDDD"),
        )
        alt_fill = PatternFill("solid", fgColor="F0F0F8")

        total_rows = 0
        for sheet_name, sheet_data in sheets_data.items():
            ws = wb.create_sheet(title=sheet_name[:31])  # Excel limit: 31 chars
            headers = sheet_data.get("headers", [])
            rows = sheet_data.get("rows", [])

            # Header row
            ws.row_dimensions[1].height = 30
            for col_idx, header in enumerate(headers, start=1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = center_align
                cell.border = thin_border

            # Data rows
            for row_idx, row in enumerate(rows, start=2):
                ws.row_dimensions[row_idx].height = 18
                fill = alt_fill if row_idx % 2 == 0 else None
                for col_idx, value in enumerate(row, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.alignment = left_align
                    cell.border = thin_border
                    if fill:
                        cell.fill = fill

            # Auto-fit column widths
            for col_idx, col_cells in enumerate(ws.columns, start=1):
                max_len = max(
                    (len(str(c.value)) if c.value is not None else 0 for c in col_cells),
                    default=8,
                )
                ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 4, 50)

            # Freeze top row
            ws.freeze_panes = "A2"
            total_rows += len(rows)

        if title:
            wb.properties.title = title

        wb.save(str(out_path))
        size_kb = out_path.stat().st_size // 1024

        sheet_names = list(sheets_data.keys())
        rel_path = _workspace_rel_path(out_path)
        return (
            f"✅ XLSX created!\n"
            f"• File: `{out_path.name}`\n"
            f"• Workspace path: `{rel_path}`\n"
            f"• Path: `{out_path}`\n"
            f"• Open: `/open {rel_path}`\n"
            f"• Size: {size_kb} KB\n"
            f"• Sheets: {', '.join(sheet_names)}\n"
            f"• Total rows: {total_rows}"
        )


# ─── DOCX ────────────────────────────────────────────────────────────────────

class DocumentCreateDocxTool(BaseTool):

    @property
    def name(self) -> str:
        return "document_create_docx"

    @property
    def description(self) -> str:
        return (
            "Create a Word document (.docx) and save it to the workspace artifacts folder. "
            "Accepts a list of sections as JSON: "
            "[{\"heading\": \"...\", \"level\": 1, \"text\": \"...\", \"bullets\": [\"...\"], \"table\": {\"headers\": [...], \"rows\": [[...]]}}, ...]. "
            "Returns the absolute file path for use with email attachments."
        )

    @property
    def category(self) -> str:
        return "DOCUMENT_TOOLS"

    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Output filename, e.g. 'report.docx'",
                },
                "title": {
                    "type": "string",
                    "description": "Document title shown as the top heading",
                },
                "sections": {
                    "type": "string",
                    "description": (
                        "JSON array of sections. Each section (all fields optional): "
                        "{\"heading\": str, \"level\": int (1-3, default 1), "
                        "\"text\": str, \"bullets\": [str], "
                        "\"table\": {\"headers\": [str], \"rows\": [[value]]}, \"image\": str (absolute path, optional)}. "
                        "Example: [{\"heading\": \"Summary\", \"level\": 1, \"text\": \"Overview text.\", \"bullets\": [\"Item A\"], \"image\": \"/path/to/chart.png\"}]"
                    ),
                },
                "author": {
                    "type": "string",
                    "description": "Author name stored in document metadata (optional)",
                },
            },
            "required": ["filename", "title", "sections"],
        }

    def execute(self, filename: str, title: str, sections: str, author: str = "") -> str:
        self._emit(f"📝 Creating DOCX: '{filename}'...")
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor as DocxRGB, Inches as DocxInches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
        except ImportError:
            return "❌ python-docx not installed. Run: pip install python-docx"

        try:
            section_data: List[dict] = json.loads(sections)
        except json.JSONDecodeError as e:
            return f"❌ Invalid JSON in 'sections': {e}"

        artifacts_dir = _get_artifacts_dir(self.scratchpad)
        out_path = artifacts_dir / _safe_filename(filename)
        if not out_path.suffix:
            out_path = out_path.with_suffix(".docx")

        doc = Document()

        # Document metadata
        core = doc.core_properties
        core.title = title
        if author:
            core.author = author

        # Heading styles tuning
        def _set_heading_color(paragraph, hex_color: str = "1A1A2E"):
            for run in paragraph.runs:
                run.font.color.rgb = DocxRGB(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16),
                )

        # Document title
        title_para = doc.add_heading(title, level=0)
        _set_heading_color(title_para, "1A1A2E")

        if author:
            para = doc.add_paragraph()
            run = para.add_run(f"Author: {author}")
            run.italic = True
            run.font.size = Pt(10)
            run.font.color.rgb = DocxRGB(0x66, 0x66, 0x66)

        doc.add_paragraph()  # spacer

        for sec in section_data:
            level = max(1, min(3, int(sec.get("level", 1))))
            heading = sec.get("heading", "")
            text = sec.get("text", "")
            bullets = sec.get("bullets", [])
            table_data = sec.get("table")

            if heading:
                h = doc.add_heading(heading, level=level)
                _set_heading_color(h, "16213E")

            if text:
                doc.add_paragraph(text)

            for bullet in bullets:
                doc.add_paragraph(bullet, style="List Bullet")

            if table_data and isinstance(table_data, dict):
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers:
                    num_cols = len(headers)
                    tbl = doc.add_table(rows=1 + len(rows), cols=num_cols)
                    tbl.style = "Table Grid"
                    # Header row
                    header_row = tbl.rows[0]
                    for col_idx, h_text in enumerate(headers):
                        cell = header_row.cells[col_idx]
                        cell.text = str(h_text)
                        run = cell.paragraphs[0].runs[0]
                        run.bold = True
                        run.font.color.rgb = DocxRGB(0xFF, 0xFF, 0xFF)
                        # Background color for header cell
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        shd = OxmlElement("w:shd")
                        shd.set(qn("w:val"), "clear")
                        shd.set(qn("w:color"), "auto")
                        shd.set(qn("w:fill"), "1A1A2E")
                        tcPr.append(shd)
                    # Data rows
                    for row_idx, row in enumerate(rows):
                        tbl_row = tbl.rows[row_idx + 1]
                        for col_idx, val in enumerate(row):
                            if col_idx < num_cols:
                                tbl_row.cells[col_idx].text = str(val)

            image_path = sec.get("image")
            if image_path and Path(image_path).exists():
                try:
                    doc.add_picture(image_path, width=DocxInches(5.0))
                except Exception as e:
                    doc.add_paragraph(f"[Image Error: {e}]")

            doc.add_paragraph()  # spacer between sections

        doc.save(str(out_path))
        size_kb = out_path.stat().st_size // 1024
        rel_path = _workspace_rel_path(out_path)

        return (
            f"✅ DOCX created!\n"
            f"• File: `{out_path.name}`\n"
            f"• Workspace path: `{rel_path}`\n"
            f"• Path: `{out_path}`\n"
            f"• Open: `/open {rel_path}`\n"
            f"• Size: {size_kb} KB\n"
            f"• Sections: {len(section_data)}"
        )
